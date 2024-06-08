from pathlib import Path
import asyncio
import os
import re
import uuid
import random
from pptx import Presentation
import json
from nonebot.log import logger
import httpx
from .openai_func import request_one
from .config import Config # type: ignore
import nonebot
from nonebot_plugin_localstore import get_cache_dir
plugin_config = Config.parse_obj(nonebot.get_driver().config.dict())


                


async def send_request(content):

    message = f"""Create a partial outline for a slide presentation with
    =====
    {content}
    ===== 
    as the outline, where the number of slides in the presentation is determined by the content.The content is in Chinese.
    You are allowed to use the following slide types:
    Title Slide - (Title, Subtitle)
    Content Slide - (Title, Content)
    Image Slide - (Title, Content, Image)
    Thanks Slide - (Title)

    Put this tag before the Title Slide: [L_TS]
    Put this tag before the Content Slide: [L_CS]
    Put this tag before the Image Slide: [L_IS]
    Put this tag before the Thanks Slide: [L_THS]
    
    Put this tag before the Title: [TITLE]
    Put this tag after the Title: [/TITLE]
    Put this tag before the Subitle: [SUBTITLE]
    Put this tag after the Subtitle: [/SUBTITLE]
    Put this tag before the Content: [CONTENT]
    Put this tag after the Content: [/CONTENT]
    Put this tag before the Image: [IMAGE]
    Put this tag after the Image: [/IMAGE]

    Put "[SLIDEBREAK]" after each slide 

    For example:
    [L_TS]
    [TITLE]Among Us[/TITLE]

    [SLIDEBREAK]

    [L_CS] 
    [TITLE]What Is Among Us?[/TITLE]
    [CONTENT]
    1. Among Us is a popular online multigptplayer game developed and published by InnerSloth.
    2. The game is set in a space-themed setting where players take on the roles of Crewmates and Impostors.
    3. The objective of Crewmates is to complete tasks and identify the Impostors among them, while the Impostors' goal is to sabotage the spaceship and eliminate the Crewmates without being caught.
    [/CONTENT]
    [IMAGE]
    Among Us game screenshot 
    [/IMAGE]

    [SLIDEBREAK]


    Make the content as detailed as possible.
    Elaborate on the content in detail, providing as much information as possibleto reduce the number of slides.
    The image content is a keyword that are highly relevant to the current content.
    REMEMBER TO PLACE a [/CONTENT] at the end of the Content.
    Do not include any special characters (?, !, ., :, ) in the Title.
    Do not include any additional information in your response and stick to the format.
"""
    messages=[
            {"role": "user", "content": message}
        ]
    data={
        "model": plugin_config.model,
        "messages": messages,
        "stream": False,
    } 
        
    try:
        response =await request_one(data)
        if response == None:
                response =await request_one(data)



    except Exception as e:
            return str(e)
    if response=="error":
        return "生成内容时发生错误，请稍后再试或尝试切换模型"
  
    response = json.loads(response.text) 
    return response['choices'][0]['message']['content']

async def generate_ppt(file_path, topic, slide_length, user_id: str):
    folder = Path("data/nonebot-plugin-multigpt/caches") / user_id
    folder.mkdir(parents=True, exist_ok=True)
    root = Presentation(file_path)


    messages=[{"role": "user", "content":f"写一个要求或主题为：\n{topic}\n的论文提纲(不包括参考文献)由{slide_length}部分组成,每个部分之间用=====加换行隔开，仅回复提纲内容"}]
    data={
        "model": plugin_config.model,
        "messages": messages,
        "stream": False,
    } 
    try:

        response =await request_one(data)
        if response == None:
                response =await request_one(data)

    except Exception as e:
        try:

            response =await request_one(data)                     
        except Exception as e:
            return str(e)
    if response=="error":
        return "生成提纲发生错误，请稍后再试或尝试切换模型"

    response = json.loads(response.text) 
    aaa=response['choices'][0]['message']['content']
    contents=aaa.split("=====")


    tasks = [send_request(content) for content in contents]
    try:
        results = await asyncio.gather(*tasks)
    except Exception as a:
        logger.info(a)
    res = "".join(results)


    def getimage(key):
        url = "https://google.serper.dev/images"

        payload = json.dumps({
            "q": key,
            "gl": "cn",
            "hl": "zh-cn"
        }, ensure_ascii=False).encode('utf-8')

        headers = {
            'X-API-KEY': plugin_config.google_key,
            'Content-Type': 'application/json; charset=utf-8'
        }

        response = httpx.post( url, headers=headers, data=payload)

        if response.status_code == 200:
            data = response.json()

            if 'images' in data and len(data['images']) > 0:
                for attempt in range(3):  # 最多尝试三次
                    random_index = random.randint(0, 4)
                    first_image = data['images'][random_index]

                    if 'imageUrl' in first_image:
                        image_url = first_image['imageUrl']
                        file_name = str(uuid.uuid4()) + ".jpg"
                        cache_dir = get_cache_dir("fetch_image")
                        file_path = os.path.join(cache_dir, file_name)
                        response = httpx.get(image_url)
                        if response.status_code == 200:
                            with open(file_path, 'wb') as f:
                                f.write(response.content)
                            # 检查文件是否存在，如果存在则返回文件路径
                            if os.path.isfile(file_path):
                                return file_path
        # 如果三次尝试都失败，返回error
        return "error"
    def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]   

    def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1 and slide.placeholders[1] is not None:
            slide.placeholders[1].text = subtitle

    def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title

    def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1 and slide.placeholders[1] is not None:
            slide.placeholders[1].text = content

    def create_title_and_content_and_image_slide(title, content, image_query):
        layout = root.slide_layouts[8] # Adjust layout as needed
        slide = root.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 2:
            slide.placeholders[2].text = content
        
        
        try:
            img_path=getimage(image_query)
            if img_path!="error":
                slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                            slide.placeholders[1].width, slide.placeholders[1].height)
        except Exception as e:
            logger.info(str(e))
            try:
                if img_path!="error":
                    img_path=getimage(image_query)
                    slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                                slide.placeholders[1].width, slide.placeholders[1].height)
            except Exception as e:
                logger.info(str(e))

    def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        for i, slide in enumerate(list_of_slides):
            slide_type = search_for_slide_type(slide)
            if slide_type == "[L_TS]":
                if i ==0:
                    create_title_slide(find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]"), 
                                    find_text_in_between_tags(slide, "[SUBTITLE]", "[/SUBTITLE]"))
            elif slide_type == "[L_CS]":
                create_title_and_content_slide(find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]"), 
                                            find_text_in_between_tags(slide, "[CONTENT]", "[/CONTENT]"))
            elif slide_type == "[L_IS]":
                create_title_and_content_and_image_slide(find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]"), 
                                                        find_text_in_between_tags(slide, "[CONTENT]", "[/CONTENT]"), 
                                                        find_text_in_between_tags(slide, "[IMAGE]", "[/IMAGE]"))
            elif slide_type == "[L_THS]":
                if i >= len(list_of_slides) - 3:  # Check if it's the last slide
                    create_section_header_slide(find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]"))

    delete_all_slides()
    logger.info(res)
    parse_response(res)

    root.save(str(folder) + f"/{topic + user_id}.pptx")

    # 删除临时文件
    dir_path = str(folder)
    prefix = "prefix_"

    for file_name in os.listdir(dir_path):
        if file_name.startswith(prefix):
            file_path = os.path.join(dir_path, file_name)
            if os.path.isfile(file_path):
                os.remove(file_path)

    return os.getcwd() + "/" + str(folder) + f"/{topic + user_id}.pptx"
